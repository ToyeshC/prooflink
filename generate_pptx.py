"""
Proof of Talent — PPTX generator
Dark Premium style matching the HTML deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.oxml.ns as pptx_ns
from lxml import etree

# ── Widescreen 16:9 ──────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Colour palette ────────────────────────────────────────────
BG       = RGBColor(0x0B, 0x0C, 0x14)
BG_ALT   = RGBColor(0x0F, 0x10, 0x19)
BG_CARD  = RGBColor(0x14, 0x15, 0x20)
BLUE     = RGBColor(0x5B, 0x8E, 0xF8)
TEAL     = RGBColor(0x00, 0xC8, 0x96)
RED      = RGBColor(0xF1, 0x5B, 0x5B)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
WHITE    = RGBColor(0xEE, 0xF0, 0xFF)
GREY     = RGBColor(0x8A, 0x8E, 0xAA)
DARK_GREY= RGBColor(0x4A, 0x4E, 0x62)
BORDER   = RGBColor(0x22, 0x24, 0x36)

# ── Helpers ───────────────────────────────────────────────────
def inch(x): return Inches(x)
def pt(x):   return Pt(x)

def fill_bg(slide, colour=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour

def add_rect(slide, x, y, w, h, fill_colour=None, line_colour=None, line_width_pt=0.75):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.RECTANGLE if False else 1,  # freeform placeholder
        inch(x), inch(y), inch(w), inch(h)
    )
    shape.line.width = pt(line_width_pt) if line_colour else 0
    if line_colour:
        shape.line.color.rgb = line_colour
    if fill_colour:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_colour
    else:
        shape.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, colour=WHITE,
             align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txb = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = pt(font_size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = colour
    run.font.name      = font_name
    return txb

def add_label(slide, text, x, y, w=4):
    """Small uppercase section label"""
    add_text(slide, text.upper(), x, y, w, 0.25,
             font_size=9, bold=True, colour=DARK_GREY, font_name="Calibri")

def add_divider(slide, x, y, w=0.5):
    """Thin blue horizontal rule"""
    d = slide.shapes.add_shape(1, inch(x), inch(y), inch(w), inch(0.03))
    d.fill.solid()
    d.fill.fore_color.rgb = BLUE
    d.line.fill.background()

def add_card(slide, x, y, w, h, fill=BG_CARD, border=BORDER):
    shape = slide.shapes.add_shape(
        1, inch(x), inch(y), inch(w), inch(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = pt(0.75)
    # rounded corners via XML
    sp = shape._element
    qn = pptx_ns.qn
    prstGeom = sp.find('.//' + qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        avLst.clear()
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 20000')
    return shape

def card_with_text(slide, cx, cy, cw, ch, label, title, body,
                   label_col=GREY, title_col=WHITE, body_col=GREY,
                   card_fill=BG_CARD, card_border=BORDER):
    add_card(slide, cx, cy, cw, ch, fill=card_fill, border=card_border)
    yy = cy + 0.13
    if label:
        add_text(slide, label.upper(), cx+0.18, yy, cw-0.36, 0.22,
                 font_size=8, bold=True, colour=label_col)
        yy += 0.22
    if title:
        add_text(slide, title, cx+0.18, yy, cw-0.36, 0.35,
                 font_size=14, bold=True, colour=title_col)
        yy += 0.32
    if body:
        add_text(slide, body, cx+0.18, yy, cw-0.36, ch - (yy - cy) - 0.1,
                 font_size=11, colour=body_col)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
fill_bg(s1, BG)

# Top badges
add_card(s1, 4.5, 0.5, 1.8, 0.35, fill=BG_CARD, border=BORDER)
add_text(s1, "Solana · SAS · x402", 4.5, 0.52, 1.8, 0.3,
         font_size=9, bold=True, colour=GREY, align=PP_ALIGN.CENTER)
add_card(s1, 6.5, 0.5, 1.2, 0.35, fill=RGBColor(0x00, 0x28, 0x1E), border=TEAL)
add_text(s1, "LIVE MVP", 6.5, 0.52, 1.2, 0.3,
         font_size=9, bold=True, colour=TEAL, align=PP_ALIGN.CENTER)

# Hero title
add_text(s1, "Proof of Talent", 1.8, 1.05, 9.7, 1.1,
         font_size=64, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_text(s1, "The decentralized oracle for verifiable human capital",
         2.0, 2.1, 9.3, 0.5,
         font_size=22, colour=GREY, align=PP_ALIGN.CENTER)

# Tagline
add_text(s1, 'AI audits developer evidence → SAS anchors skill attestations → agents pay $0.001 per lookup',
         2.2, 2.65, 8.9, 0.4,
         font_size=14, colour=DARK_GREY, align=PP_ALIGN.CENTER)

# Demo link box
add_card(s1, 4.3, 3.15, 4.7, 0.45, fill=RGBColor(0x0D, 0x1A, 0x38), border=BLUE)
add_text(s1, "▶  Live Demo — proof-of-talent.vercel.app",
         4.3, 3.17, 4.7, 0.4,
         font_size=13, bold=True, colour=BLUE, align=PP_ALIGN.CENTER)

# Pipeline — 4 nodes + connecting line
pipe_y   = 3.85
node_w   = 2.2
node_h   = 1.2
gap      = 0.4
start_x  = 0.95
nodes = [
    ("🎓", "Evidence", "GitHub + Canvas\nZIP exports"),
    ("🧠", "AI Oracle", "Skill · Confidence\nEvidence summary"),
    ("🔐", "SAS + ZK", "On-chain attestation\nZK-privacy layer"),
    ("⚡", "x402 API", "Machine-readable\nmicropayment-gated"),
]
# connector line
line = s1.shapes.add_shape(1,
    inch(start_x + node_w/2), inch(pipe_y + 0.28),
    inch((node_w + gap) * 3), inch(0.03))
line.fill.solid(); line.fill.fore_color.rgb = BLUE
line.line.fill.background()

for i, (icon, label, sub) in enumerate(nodes):
    nx = start_x + i * (node_w + gap)
    border_col = TEAL if i == 3 else BLUE if i in (1,2) else BORDER
    add_card(s1, nx, pipe_y, node_w, node_h, fill=BG_CARD, border=border_col)
    add_text(s1, icon,  nx, pipe_y + 0.08, node_w, 0.4, font_size=24, align=PP_ALIGN.CENTER)
    add_text(s1, label, nx, pipe_y + 0.52, node_w, 0.3,
             font_size=13, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    add_text(s1, sub,   nx, pipe_y + 0.80, node_w, 0.38,
             font_size=10, colour=GREY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ═══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s2, BG_ALT)

add_label(s2, "The Problem", 0.65, 0.45)
add_text(s2, "Academic signals are noisy, unstandardised\nproxies for actual skill.",
         0.65, 0.7, 8.5, 1.1,
         font_size=36, bold=True, colour=WHITE)
add_divider(s2, 0.65, 1.82)

# Left card — problem bullets
add_card(s2, 0.65, 2.0, 5.5, 4.6, fill=RGBColor(0x1A, 0x10, 0x10), border=RED)
add_text(s2, "THE VALIDITY PROBLEM TODAY", 0.83, 2.15, 5.1, 0.25,
         font_size=9, bold=True, colour=RED)
bullets = [
    "Degrees measure compliance — not competence",
    "Credentials locked in institutional silos",
    "Manual verification: $40–$200 per check",
    "AI agents can't query unstructured claims",
    "No cheap API-native trust primitive exists",
]
for i, b in enumerate(bullets):
    add_text(s2, f"•  {b}", 0.83, 2.48 + i * 0.72, 5.1, 0.65,
             font_size=14, colour=RGBColor(0xCC, 0xBB, 0xBB))

# Right — stat cards
add_card(s2, 6.5, 2.0, 6.18, 2.1, fill=RGBColor(0x0D, 0x12, 0x28), border=BLUE)
add_text(s2, "THE COST GAP", 6.68, 2.15, 5.8, 0.25, font_size=9, bold=True, colour=BLUE)
add_text(s2, "$40–$200", 6.68, 2.45, 2.6, 0.7, font_size=40, bold=True, colour=RED)
add_text(s2, "cost per background check\n(Checkr / Sterling / HireRight)",
         6.68, 3.05, 2.8, 0.55, font_size=11, colour=GREY)
# divider
d = s2.shapes.add_shape(1, inch(9.5), inch(2.45), inch(0.03), inch(1.2))
d.fill.solid(); d.fill.fore_color.rgb = BORDER; d.line.fill.background()
add_text(s2, "$0.001", 9.7, 2.45, 2.6, 0.7, font_size=40, bold=True, colour=TEAL)
add_text(s2, "our API lookup fee\n75% gross margin",
         9.7, 3.05, 2.6, 0.55, font_size=11, colour=GREY)

add_card(s2, 6.5, 4.3, 6.18, 2.3, fill=BG_CARD, border=BORDER)
add_text(s2, "THE MACHINE SCREENING PROBLEM", 6.68, 4.45, 5.8, 0.25,
         font_size=9, bold=True, colour=GREY)
add_text(s2,
    "AI recruiting agents evaluate 100,000 candidates overnight — "
    "but cannot trust any claim they read. There is no structured, "
    "machine-readable skill primitive they can query programmatically.",
    6.68, 4.78, 5.82, 1.6, font_size=13, colour=GREY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — WHY NOW
# ═══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s3, BG)

add_label(s3, "Why Now", 0.65, 0.45)
add_text(s3, "Four forces converged. The window is this year.",
         0.65, 0.7, 11.0, 0.85, font_size=36, bold=True, colour=WHITE)
add_divider(s3, 0.65, 1.58)

why_items = [
    ("01", "AI agents now handle hiring at scale",
     "Automated screening went from 'nice to have' to default in 2025. "
     "Agent-to-agent hiring is real. They need trust primitives machines can query."),
    ("02", "Solana Attestation Service launched",
     "SAS provides a native, interoperable on-chain credential layer with composable schemas — "
     "this infrastructure didn't exist 18 months ago."),
    ("03", "Micropayments at machine scale are viable",
     "Solana's parallel tx engine makes $0.001-per-lookup viable at 75% gross margin. "
     "x402 brings HTTP-native monetisation to APIs."),
    ("04", "Public artifacts are the primary signal",
     "GitHub history and Canvas course submissions are the most honest technical evidence available — "
     "queryable, unfiltered, student-initiated."),
]

for i, (num, title, body) in enumerate(why_items):
    col = i % 2
    row = i // 2
    cx = 0.65 + col * 6.35
    cy = 1.8  + row * 2.55
    add_card(s3, cx, cy, 6.0, 2.35, fill=BG_CARD, border=BORDER)
    add_text(s3, num, cx + 0.2, cy + 0.15, 0.8, 0.6,
             font_size=28, bold=True, colour=BLUE)
    add_text(s3, title, cx + 0.2, cy + 0.65, 5.6, 0.4,
             font_size=14, bold=True, colour=WHITE)
    add_text(s3, body, cx + 0.2, cy + 1.05, 5.6, 1.1,
             font_size=11, colour=GREY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION
# ═══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s4, BG_ALT)

add_label(s4, "The Solution", 0.65, 0.45)
add_text(s4, "Raw artifacts → sovereign, verifiable,\nmachine-readable skill proof.",
         0.65, 0.7, 11.0, 1.1, font_size=36, bold=True, colour=WHITE)
add_divider(s4, 0.65, 1.82)

sol_items = [
    ("01 / INGEST", "Multi-source evidence",
     "GitHub repos, PRs, commit history. Canvas ZIP course exports. "
     "Any structured artifact a student initiates — not what an institution certifies.",
     BLUE, BORDER),
    ("02 / INFER", "AI skill oracle",
     "LLM interprets raw evidence. Outputs skill slug, confidence score, evidence summary. "
     "Bypasses institutional single points of failure.",
     BLUE, BORDER),
    ("03 / ANCHOR", "SAS + ZK-privacy",
     "Claims written via Solana Attestation Service. ZK-privacy layer protects raw evidence. "
     "Optimistic dispute mechanism with staked proofs for re-evaluation.",
     TEAL, RGBColor(0x00, 0x40, 0x30)),
    ("04 / SERVE", "x402 verification API",
     "Profile lookup (free). Basic verify $0.001. Premium verify $0.005 with ZK proof. "
     "Payment-gated via x402 stablecoin — native to AI agents.",
     TEAL, RGBColor(0x00, 0x40, 0x30)),
]

card_w = 2.9
for i, (label, title, body, lbl_col, bdr) in enumerate(sol_items):
    cx = 0.65 + i * (card_w + 0.22)
    add_card(s4, cx, 2.0, card_w, 3.6, fill=BG_CARD, border=bdr)
    # top accent bar
    bar = s4.shapes.add_shape(1, inch(cx), inch(2.0), inch(card_w), inch(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = lbl_col; bar.line.fill.background()
    add_text(s4, label, cx+0.15, 2.12, card_w-0.3, 0.25,
             font_size=9, bold=True, colour=lbl_col)
    add_text(s4, title, cx+0.15, 2.38, card_w-0.3, 0.45,
             font_size=15, bold=True, colour=WHITE)
    add_text(s4, body,  cx+0.15, 2.85, card_w-0.3, 2.55,
             font_size=11, colour=GREY)

# JSON code block
add_card(s4, 0.65, 5.72, 12.0, 1.45, fill=RGBColor(0x07, 0x08, 0x0F), border=DARK_GREY)
add_text(s4,
    '{ "skill": "rust",  "confidence": 0.91,  "evidence": "14 Rust repos · Anchor framework · 847 commits",\n'
    '  "sasAttestation": "3xY8...k9mP",  "zkProof": "groth16:Ax7B...",  "sovereign": true  // wallet-bound, portable }',
    0.82, 5.82, 11.7, 1.2,
    font_size=11, colour=RGBColor(0x7B, 0x8E, 0xC8), font_name="Courier New")

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — USER FLOW
# ═══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s5, BG)

add_label(s5, "How It Works", 0.65, 0.45)
add_text(s5, "One endpoint for breadth. Paid depth on demand.",
         0.65, 0.7, 11.0, 0.7, font_size=36, bold=True, colour=WHITE)
add_divider(s5, 0.65, 1.43)

lanes = [
    {
        "title": "👤  Student / Developer",
        "colour": RGBColor(0x3A, 0x5E, 0xC8),
        "steps": [
            ("Submit GitHub + wallet address", BLUE, RGBColor(0x0D, 0x1A, 0x38)),
            ("Oracle processes in background…", DARK_GREY, BG_CARD),
            ("Receive SAS attestation + explorer link", BLUE, RGBColor(0x0D, 0x1A, 0x38)),
            ("Profile: sovereign, portable, queryable", DARK_GREY, BG_CARD),
            ("Dispute: stake SOL to challenge a claim", AMBER, RGBColor(0x1A, 0x14, 0x05)),
        ],
    },
    {
        "title": "⚙️  Proof of Talent Oracle",
        "colour": BLUE,
        "steps": [
            ("Fetch GitHub repos · PRs · commits", BLUE, RGBColor(0x0D, 0x1A, 0x38)),
            ("Fetch Canvas ZIP exports (if provided)", BLUE, RGBColor(0x0D, 0x1A, 0x38)),
            ("LLM scores skill · confidence · evidence", BLUE, RGBColor(0x0D, 0x1A, 0x38)),
            ("Write SAS attestation to Solana", BLUE, RGBColor(0x0D, 0x1A, 0x38)),
            ("ZK-proof generated · raw data private", BLUE, RGBColor(0x0D, 0x1A, 0x38)),
            ("Optimistic dispute: re-evaluate on stake", AMBER, RGBColor(0x1A, 0x14, 0x05)),
        ],
    },
    {
        "title": "🤖  Recruiter / AI Agent",
        "colour": TEAL,
        "steps": [
            ("GET /api/profile/:wallet → skill map", TEAL, RGBColor(0x05, 0x20, 0x18)),
            ("Free. Filter candidates at breadth.", DARK_GREY, BG_CARD),
            ("GET /api/verify → 402 Payment Required", RED, RGBColor(0x20, 0x08, 0x08)),
            ("x402: $0.001 basic · $0.005 ZK-backed", DARK_GREY, BG_CARD),
            ("verified: true + evidence + SAS tx link", TEAL, RGBColor(0x05, 0x20, 0x18)),
        ],
    },
]

lane_w = 3.95
for li, lane in enumerate(lanes):
    lx = 0.65 + li * (lane_w + 0.22)
    # header
    add_card(s5, lx, 1.6, lane_w, 0.42, fill=RGBColor(0x12, 0x14, 0x24), border=lane["colour"])
    add_text(s5, lane["title"], lx+0.12, 1.65, lane_w-0.24, 0.32,
             font_size=11, bold=True, colour=lane["colour"])
    # steps
    sy = 2.1
    for txt, col, fill in lane["steps"]:
        add_card(s5, lx, sy, lane_w, 0.72, fill=fill, border=col)
        add_text(s5, txt, lx+0.14, sy+0.12, lane_w-0.28, 0.5,
                 font_size=11, colour=col if col != DARK_GREY else GREY)
        sy += 0.8

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s6, BG_ALT)

add_label(s6, "Business Model", 0.65, 0.45)
add_text(s6, "Infrastructure economics: revenue scales with AI adoption.",
         0.65, 0.7, 11.5, 0.7, font_size=34, bold=True, colour=WHITE)
add_divider(s6, 0.65, 1.43)

# Left — 4 revenue streams
rev_streams = [
    ("Basic Verify · x402",    "$0.001 /call",   "Per-call stablecoin · 75% gross margin · $0.00025 tx cost · 30,000× cheaper than Checkr", TEAL),
    ("Premium Verify · ZK",    "$0.005 /call",   "Full evidence chain + ZK proof + dispute guarantee · 95% gross margin", TEAL),
    ("Recruiter SaaS Plans",   "$499–$9,999/mo", "Team quotas · ATS integrations · analytics · enterprise custom schemas", BLUE),
    ("Issuer Plans + Disputes","$999/mo + %",    "Bootcamps/universities issue attestations · 10% protocol cut on challenged claims", AMBER),
]

for i, (name, price, desc, col) in enumerate(rev_streams):
    ry = 1.62 + i * 1.3
    add_card(s6, 0.65, ry, 5.8, 1.18, fill=BG_CARD, border=BORDER)
    bar = s6.shapes.add_shape(1, inch(0.65), inch(ry), inch(0.05), inch(1.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    add_text(s6, name,  0.82, ry+0.1,  3.8, 0.32, font_size=14, bold=True, colour=WHITE)
    add_text(s6, price, 4.7,  ry+0.1,  1.6, 0.32, font_size=16, bold=True, colour=col,
             align=PP_ALIGN.RIGHT)
    add_text(s6, desc,  0.82, ry+0.48, 5.4, 0.55, font_size=11, colour=GREY)

# Right — Revenue scale table
add_card(s6, 6.8, 1.62, 5.88, 5.5, fill=RGBColor(0x0D, 0x1A, 0x38), border=BLUE)
add_text(s6, "Revenue at Machine Scale", 7.0, 1.75, 5.5, 0.38,
         font_size=16, bold=True, colour=WHITE)
add_text(s6, "@ $0.001 base · 75% gross margin",
         7.0, 2.1, 5.5, 0.26, font_size=9, colour=GREY)

# Table header
header_y = 2.48
for xi, (hdr, hx, hw) in enumerate([
    ("Verifications/Day", 6.9, 2.3),
    ("API Revenue",       9.3, 1.5),
    ("+ SaaS",           10.85, 1.0),
    ("ARR",              11.92, 0.65),
]):
    add_text(s6, hdr, hx, header_y, hw, 0.28,
             font_size=9, bold=True, colour=DARK_GREY,
             align=PP_ALIGN.RIGHT if xi > 0 else PP_ALIGN.LEFT)

# Divider under header
d = s6.shapes.add_shape(1, inch(6.9), inch(2.78), inch(5.6), inch(0.02))
d.fill.solid(); d.fill.fore_color.rgb = BORDER; d.line.fill.background()

rows = [
    ("1,000",   "$365K",   "$240K",  "$605K",  False),
    ("10,000",  "$3.65M",  "$960K",  "~$4.6M", True),
    ("100,000", "$36.5M",  "$3.6M",  "~$40M",  False),
]
for ri, (vol, api, saas, arr, highlight) in enumerate(rows):
    ry2 = 2.9 + ri * 1.28
    if highlight:
        add_card(s6, 6.9, ry2-0.08, 5.58, 1.15,
                 fill=RGBColor(0x0D, 0x22, 0x48), border=BLUE)
    col_arr  = TEAL if highlight else WHITE
    col_rest = WHITE if highlight else GREY
    add_text(s6, vol,  6.9,  ry2, 2.2, 0.5, font_size=16, bold=highlight, colour=col_rest)
    add_text(s6, api,  9.2,  ry2, 1.6, 0.5, font_size=16, bold=highlight, colour=col_rest,
             align=PP_ALIGN.RIGHT)
    add_text(s6, saas, 10.84,ry2, 1.0, 0.5, font_size=14, colour=col_rest,
             align=PP_ALIGN.RIGHT)
    add_text(s6, arr,  11.9, ry2, 0.68,0.5, font_size=16, bold=True, colour=col_arr,
             align=PP_ALIGN.RIGHT)
    if not highlight:
        d2 = s6.shapes.add_shape(1, inch(6.9), inch(ry2+0.62), inch(5.58), inch(0.015))
        d2.fill.solid(); d2.fill.fore_color.rgb = BORDER; d2.line.fill.background()

# Footer note
add_text(s6, "▲ Gross margin: 75% basic · 95% premium     ▲ Checkr: $29.99/check — we're 30,000× cheaper",
         6.9, 6.75, 5.6, 0.28, font_size=9, colour=DARK_GREY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — STATUS + ROADMAP
# ═══════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(s7, BG)

add_label(s7, "Product Status", 0.65, 0.45)
add_text(s7, "Core loop works today. Roadmap upgrades trust guarantees.",
         0.65, 0.7, 11.5, 0.7, font_size=34, bold=True, colour=WHITE)
add_divider(s7, 0.65, 1.43)

# Live MVP badge
add_card(s7, 0.65, 1.62, 1.6, 0.36, fill=RGBColor(0x00, 0x28, 0x1E), border=TEAL)
add_text(s7, "● LIVE MVP", 0.65, 1.65, 1.6, 0.3,
         font_size=10, bold=True, colour=TEAL, align=PP_ALIGN.CENTER)

live_items = [
    "GitHub evidence ingestion + parsing",
    "AI skill inference (Anthropic / OpenRouter)",
    "On-chain proof anchor (Solana memo tx)",
    "Profile + verify API endpoints live",
    "402 payment-gated flow in demo UI",
]
for i, item in enumerate(live_items):
    iy = 2.1 + i * 0.95
    add_card(s7, 0.65, iy, 5.8, 0.82,
             fill=BG_CARD, border=RGBColor(0x00, 0x40, 0x30))
    # left bar
    bar = s7.shapes.add_shape(1, inch(0.65), inch(iy), inch(0.05), inch(0.82))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    add_text(s7, "✓", 0.8, iy+0.2, 0.35, 0.4, font_size=16, bold=True, colour=TEAL)
    add_text(s7, item, 1.2, iy+0.18, 5.1, 0.45, font_size=13, colour=WHITE)

# Roadmap badge
add_card(s7, 6.88, 1.62, 2.0, 0.36, fill=RGBColor(0x0D, 0x1A, 0x38), border=BLUE)
add_text(s7, "→ NEXT MILESTONES", 6.88, 1.65, 2.0, 0.3,
         font_size=10, bold=True, colour=BLUE, align=PP_ALIGN.CENTER)

next_items = [
    ("Full SAS issuance lifecycle",
     "Schema · Credential · Attestation — fully interoperable"),
    ("Canvas ZIP ingestion pipeline",
     "Course submissions · grades · institutional artifacts"),
    ("ZK-privacy layer + dispute mechanism",
     "Raw evidence protected · staked proofs for re-evaluation"),
    ("Production indexer + hardened x402",
     "Postgres · chain indexer · amount/recipient validation"),
]
for i, (title, sub) in enumerate(next_items):
    ny = 2.1 + i * 1.22
    add_card(s7, 6.88, ny, 5.8, 1.05,
             fill=BG_CARD, border=RGBColor(0x1A, 0x30, 0x5C))
    bar2 = s7.shapes.add_shape(1, inch(6.88), inch(ny), inch(0.05), inch(1.05))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = BLUE; bar2.line.fill.background()
    add_text(s7, title, 7.05, ny+0.12, 5.5, 0.36, font_size=14, bold=True, colour=WHITE)
    add_text(s7, sub,   7.05, ny+0.52, 5.5, 0.45, font_size=11, colour=GREY)

# ── Save ──────────────────────────────────────────────────────
OUT = "/Users/toyesh/Documents/Personal/Coding/my-solana-agent/proof-of-talent-pitch.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
